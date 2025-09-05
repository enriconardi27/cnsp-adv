import streamlit as st
import os
import json
import shutil
import csv
import pandas as pd
import docx
import zipfile
import py7zr
import requests
import bcrypt
from bs4 import BeautifulSoup
from pptx import Presentation
from datetime import datetime
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain.vectorstores import FAISS
from langchain.chains.conversational_retrieval.base import ConversationalRetrievalChain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.document_loaders import PyPDFLoader
from langchain.docstore.document import Document
from langchain.prompts import PromptTemplate

try:
    GOOGLE_API_KEY = st.secrets.get("GOOGLE_API_KEY", "INSERISCI_LA_TUA_API_KEY_QUI")
except (AttributeError, FileNotFoundError):
    GOOGLE_API_KEY = "INSERISCI_LA_TUA_API_KEY_QUI"

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
VECTOR_STORE_PATH = os.path.join(BASE_DIR, "vector_store")
USER_DATA_PATH = os.path.join(BASE_DIR, "user_data")
TEMP_FILES_PATH = os.path.join(BASE_DIR, "temp")
FEEDBACK_FILE_PATH = os.path.join(BASE_DIR, "feedback_log.csv")
USERS_FILE_PATH = os.path.join(BASE_DIR, "users.json")

def load_users():
    if not os.path.exists(USERS_FILE_PATH):
        return {}
    with open(USERS_FILE_PATH, "r") as f:
        return json.load(f)

def save_users(users_data):
    with open(USERS_FILE_PATH, "w") as f:
        json.dump(users_data, f, indent=4)

def hash_password(password):
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(stored_password, provided_password):
    return bcrypt.checkpw(provided_password.encode('utf-8'), stored_password.encode('utf-8'))

def get_user_dir(username):
    return os.path.join(USER_DATA_PATH, username)

def get_user_chat_file_path(username, chat_id):
    user_dir = get_user_dir(username)
    return os.path.join(user_dir, "chats", f"{chat_id}.json")

def load_user_data(username):
    user_dir = get_user_dir(username)
    chats_dir = os.path.join(user_dir, "chats")

    if not os.path.exists(chats_dir):
        os.makedirs(chats_dir)
        return {"chats": {}, "active_chat_id": None}

    user_data = {"chats": {}}
    for filename in os.listdir(chats_dir):
        if filename.endswith(".json"):
            chat_id = os.path.splitext(filename)[0]
            with open(os.path.join(chats_dir, filename), "r") as f:
                user_data["chats"][chat_id] = json.load(f)

    if user_data["chats"]:
        latest_chat_id = max(user_data["chats"].keys())
        user_data["active_chat_id"] = latest_chat_id
    else:
        user_data["active_chat_id"] = None

    return user_data

def save_active_chat(username):
    if "user_data" in st.session_state and "active_chat_id" in st.session_state.user_data:
        active_chat_id = st.session_state.user_data["active_chat_id"]
        if active_chat_id:
            active_chat_data = st.session_state.user_data["chats"][active_chat_id]
            chat_file_path = get_user_chat_file_path(username, active_chat_id)

            os.makedirs(os.path.dirname(chat_file_path), exist_ok=True)

            with open(chat_file_path, "w") as f:
                json.dump(active_chat_data, f)

def scrape_faq_data():
    base_url = "https://www.acquistinretepa.it/opencms/opencms/faq.html"
    page_number = 1
    scraped_documents = []

    st.write("Inizio importazione FAQ...")
    progress_bar = st.progress(0)

    while True:
        current_url = f"{base_url}?page={page_number}"
        try:
            response = requests.get(current_url)
            response.raise_for_status()

            soup = BeautifulSoup(response.content, 'html.parser')
            qa_pairs = soup.select("div.row.question")
            if not qa_pairs:
                break

            for pair in qa_pairs:
                question_tag = pair.find('h3')
                answer_tag = pair.find('div', class_='testo')
                if question_tag and answer_tag:
                    question = question_tag.get_text(strip=True)
                    answer = answer_tag.get_text(strip=True)
                    content = f"Domanda: {question}\nRisposta: {answer}"
                    metadata = {'source': current_url, 'question': question}
                    scraped_documents.append(Document(page_content=content, metadata=metadata))

            st.write(f"Pagina {page_number} importata con successo.")
            progress_bar.progress(min(page_number / 16, 1.0))
            page_number += 1

        except requests.RequestException as e:
            st.error(f"Errore di rete durante l'accesso alla pagina {page_number}: {e}")
            break

    progress_bar.progress(1.0)
    st.write("Importazione FAQ completata.")
    return scraped_documents

def switch_chat(chat_id):
    st.session_state.user_data["active_chat_id"] = chat_id
    active_chat = st.session_state.user_data["chats"][chat_id]
    st.session_state.messages = active_chat["messages"]
    st.session_state.chat_history_tuples = [
        (msg["content"], st.session_state.messages[i+1]["content"])
        for i, msg in enumerate(st.session_state.messages)
        if msg["role"] == "user" and i+1 < len(st.session_state.messages)
    ]

def create_new_chat():
    chat_id = f"chat_{int(datetime.now().timestamp())}"
    chat_name = f"Chat del {datetime.now().strftime('%d/%m %H:%M')}"
    st.session_state.user_data["chats"][chat_id] = {"name": chat_name, "messages": []}
    switch_chat(chat_id)

def get_chat_summary(user_message, assistant_message):
    try:
        llm = ChatGoogleGenerativeAI(model="gemini-1.5-pro-latest", google_api_key=GOOGLE_API_KEY, temperature=0.0)
        prompt = f"""
        Basandoti sulla seguente conversazione, crea un titolo molto breve (massimo 5 parole) che ne riassuma l'argomento principale.

        CONVERSAZIONE:
        Utente: "{user_message}"
        Assistente: "{assistant_message}"

        TITOLO BREVE:
        """
        response = llm.invoke(prompt)
        summary = response.content.strip().replace('"', '')
        return summary
    except Exception as e:
        print(f"Error generating chat summary: {e}")
        return f"Chat del {datetime.now().strftime('%d/%m')}"

def get_suggestions(user_question, retriever):
    try:
        similar_docs = retriever.get_relevant_documents(user_question)
        if not similar_docs:
            return None

        context = "\n\n".join([doc.page_content for doc in similar_docs])
        llm = ChatGoogleGenerativeAI(model="gemini-1.5-pro-latest", google_api_key=GOOGLE_API_KEY, temperature=0.0)
        suggestion_prompt = f"""
        Un utente ha cercato "{user_question}" ma non ha trovato una risposta diretta.
        Basandoti sul seguente contesto estratto dai documenti, suggerisci 3 argomenti o termini di ricerca alternativi e pertinenti.
        Rispondi solo con una lista puntata di suggerimenti in italiano.

        Contesto:
        {context}

        Suggerimenti:
        """
        response = llm.invoke(suggestion_prompt)
        return response.content
    except Exception as e:
        print(f"Error generating suggestions: {e}")
        return None

def handle_user_input(user_question):
    if st.session_state.conversation:
        active_chat_id = st.session_state.user_data["active_chat_id"]
        if not active_chat_id:
            st.warning("Nessuna chat attiva. Per favore, crea una nuova chat.")
            return

        response = st.session_state.conversation({'question': user_question, 'chat_history': st.session_state.get('chat_history_tuples', [])})

        if not response['source_documents']:
            suggestions = get_suggestions(user_question, st.session_state.conversation.retriever)
            if suggestions:
                answer = f"Non ho trovato risultati diretti per '{user_question}'.\n\nForse cercavi:\n{suggestions}"
            else:
                answer = f"Mi dispiace, non ho trovato alcuna informazione relativa a '{user_question}' nei documenti a mia disposizione."
            citations = []
        else:
            answer = response['answer']
            st.session_state.chat_history_tuples.append((user_question, answer))

            citations = []
            for doc in response['source_documents']:
                meta = doc.metadata
                source = meta.get('source', 'N/D')
                
                details = []
                if 'page' in meta: details.append(f"Pagina: {meta['page'] + 1}")
                if 'paragraph' in meta: details.append(f"Paragrafo: {meta['paragraph']}")
                if 'row' in meta: details.append(f"Riga: {meta['row']}")
                if 'line' in meta: details.append(f"Riga: {meta['line']}") 
                if 'slide' in meta: details.append(f"Slide: {meta['slide']}")
                
                if details:
                    citation_detail = f"**{source}** ({', '.join(details)})"
                else:
                    citation_detail = f"**{source}**"
                
                if citation_detail not in citations:
                    citations.append(citation_detail)

        assistant_message = {"role": "assistant", "content": answer, "question": user_question, "citations": citations}
        st.session_state.user_data["chats"][active_chat_id]["messages"].append(assistant_message)
        st.session_state.messages = st.session_state.user_data["chats"][active_chat_id]["messages"]

        active_chat = st.session_state.user_data["chats"][active_chat_id]
        if len(active_chat["messages"]) == 2:
            new_title = get_chat_summary(user_question, answer)
            if new_title:
                st.session_state.user_data["chats"][active_chat_id]["name"] = new_title

def find_supported_files_in_path(directory_path, allowed_extensions):
    found_files = []
    for root, _, files in os.walk(directory_path):
        for file in files:
            if any(file.lower().endswith(ext) for ext in allowed_extensions):
                found_files.append(os.path.join(root, file))
    return found_files

def get_documents_with_detailed_metadata(uploaded_files):
    all_documents = []
    allowed_single_file_ext = ['.pdf', '.docx', '.csv', '.xlsx', '.pptx', '.aspx', '.txt', '.md']

    for uploaded_file in uploaded_files:
        file_path = os.path.join(TEMP_FILES_PATH, uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        file_extension = os.path.splitext(uploaded_file.name)[1].lower()

        if file_extension in ['.zip', '.7z']:
            extraction_path = os.path.join(TEMP_FILES_PATH, uploaded_file.name + "_extracted")
            if not os.path.exists(extraction_path):
                os.makedirs(extraction_path)
            
            try:
                if file_extension == '.zip':
                    with zipfile.ZipFile(file_path, 'r') as zf:
                        zf.extractall(extraction_path)
                elif file_extension == '.7z':
                    with py7zr.SevenZipFile(file_path, mode='r') as z:
                        z.extractall(path=extraction_path)

                files_to_process = find_supported_files_in_path(extraction_path, allowed_single_file_ext)
                
                if not files_to_process:
                    st.warning(f"L'archivio '{uploaded_file.name}' è vuoto o non contiene file supportati.")
                    continue

                with st.spinner(f"Processando i file dall'archivio '{uploaded_file.name}'..."):
                    for doc_path in files_to_process:
                        doc_name = os.path.basename(doc_path)
                        all_documents.extend(process_single_file(doc_path, doc_name, archive_name=uploaded_file.name))

            except Exception as e:
                st.error(f"Errore durante l'estrazione dell'archivio {uploaded_file.name}: {e}")

        elif file_extension in allowed_single_file_ext:
            all_documents.extend(process_single_file(file_path, uploaded_file.name))
        
        else:
            st.warning(f"Il tipo di file '{uploaded_file.name}' non è supportato e sarà ignorato.")

    return all_documents


def process_single_file(file_path_or_obj, file_name, archive_name=None):
    documents = []
    source_display = f"{archive_name} -> {file_name}" if archive_name else file_name
    
    file_extension = os.path.splitext(file_name)[1].lower()
    
    temp_file_path = file_path_or_obj if isinstance(file_path_or_obj, str) else os.path.join(TEMP_FILES_PATH, file_name)
    if not isinstance(file_path_or_obj, str) and hasattr(file_path_or_obj, 'getbuffer'):
        with open(temp_file_path, "wb") as f: f.write(file_path_or_obj.getbuffer())

    try:
        if file_extension == '.pdf':
            loader = PyPDFLoader(temp_file_path)
            pages = loader.load_and_split()
            for page in pages:
                page.metadata['source'] = source_display
            documents.extend(pages)
            
        elif file_extension == '.docx':
            doc = docx.Document(temp_file_path)
            for i, p in enumerate(doc.paragraphs):
                if p.text.strip():
                    documents.append(Document(page_content=p.text, metadata={'source': source_display, 'paragraph': i + 1}))
                    
        elif file_extension == '.pptx':
            prs = Presentation(temp_file_path)
            for i, slide in enumerate(prs.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"): slide_text += shape.text + "\n"
                if slide_text.strip():
                    documents.append(Document(page_content=slide_text, metadata={'source': source_display, 'slide': i + 1}))
                    
        elif file_extension == '.csv':
            with open(temp_file_path, mode='r', encoding='utf-8') as csvfile:
                reader = csv.reader(csvfile)
                header = next(reader, None)
                for i, row in enumerate(reader):
                    row_content = ", ".join(f"{header[j]}: {val}" for j, val in enumerate(row) if header)
                    documents.append(Document(page_content=row_content, metadata={'source': source_display, 'row': i + 2}))

        elif file_extension == '.xlsx':
            df = pd.read_excel(temp_file_path)
            for i, row in df.iterrows():
                row_content = ", ".join(f"{col}: {val}" for col, val in row.items() if pd.notna(val))
                documents.append(Document(page_content=row_content, metadata={'source': source_display, 'row': i + 2}))

        elif file_extension in ['.aspx', '.txt', '.md']:
            with open(temp_file_path, "r", encoding='utf-8', errors='ignore') as f:
                for i, line_text in enumerate(f):
                    if line_text.strip():
                        if file_extension == '.aspx':
                            soup = BeautifulSoup(line_text, 'html.parser')
                            clean_text = soup.get_text(strip=True)
                        else:
                            clean_text = line_text
                        
                        if clean_text:
                            documents.append(Document(page_content=clean_text, metadata={'source': source_display, 'line': i + 1}))

    except Exception as e:
        st.warning(f"Impossibile processare {file_name}: {e}")
    return documents


def get_text_chunks(documents):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    return text_splitter.split_documents(documents)

def get_vector_store(text_chunks):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=GOOGLE_API_KEY)
        vector_store = FAISS.from_documents(text_chunks, embedding=embeddings)
        vector_store.save_local(VECTOR_STORE_PATH)
        return vector_store
    except Exception as e: st.error(f"Errore creazione Vector Store: {e}"); return None

def get_conversational_chain(vector_store):
    try:
        llm = ChatGoogleGenerativeAI(model="gemini-1.5-pro-latest", google_api_key=GOOGLE_API_KEY, temperature=0.5)
        
        template = """
        Sei 'Consip Advisor', un assistente AI amichevole, professionale e colloquiale. 
        Il tuo obiettivo primario è assistere gli utenti rispondendo alle loro domande in modo chiaro e utile, basandoti sulle informazioni trovate nel contesto fornito.
        
        ISTRUZIONI:
        1. Rispondi sempre in italiano, con un tono naturale e discorsivo.
        2. Quando la risposta è presente nel contesto, sintetizzala con parole tue invece di copiarla. Inizia la risposta in modo diretto e poi, se necessario, aggiungi dettagli.
        3. Se il contesto non contiene una risposta chiara, **non inventare informazioni**. Invece, spiega gentilmente che non hai trovato una corrispondenza esatta nei documenti. Puoi offrire informazioni correlate se presenti, o suggerire come l'utente potrebbe riformulare la domanda.
        4. Usa la cronologia della chat per capire il contesto della conversazione e fornire risposte pertinenti alle domande successive.
        5. Sii sempre cortese e concludi le tue risposte in modo amichevole.

        Contesto: {context}
        Cronologia Chat: {chat_history}
        Domanda: {question}
        Risposta:
        """
        prompt = PromptTemplate(template=template, input_variables=["context", "chat_history", "question"])

        chain = ConversationalRetrievalChain.from_llm(
            llm=llm, 
            retriever=vector_store.as_retriever(search_kwargs={"k": 5}), 
            return_source_documents=True,
            combine_docs_chain_kwargs={"prompt": prompt}
        )
        return chain
    except Exception as e: st.error(f"Errore creazione catena conversazionale: {e}"); return None

def save_feedback(feedback_data):
    file_exists = os.path.exists(FEEDBACK_FILE_PATH)
    with open(FEEDBACK_FILE_PATH, 'a', newline='', encoding='utf-8') as f:
        writer = csv.DictWriter(f, fieldnames=['timestamp', 'username', 'question', 'response', 'rating', 'comment'])
        if not file_exists: writer.writeheader()
        writer.writerow(feedback_data)

def show_login_or_register():
    st.sidebar.title("Accesso Riservato")
    choice = st.sidebar.radio("Scegli un'opzione", ["Login", "Registrati"])

    if choice == "Login":
        login_page()
    else:
        register_page()
    return False

def login_page():
    st.header("Login")
    with st.form("login_form"):
        username = st.text_input("Username").lower()
        password = st.text_input("Password", type="password")
        submitted = st.form_submit_button("Accedi")

        if submitted:
            users = load_users()
            if username in users and verify_password(users[username]["password"], password):
                st.session_state.logged_in = True
                st.session_state.current_user = username
                st.success("Accesso effettuato con successo!")
                st.rerun()
            else:
                st.error("Username o password non corretti.")

def register_page():
    st.header("Crea un nuovo account")
    with st.form("register_form"):
        username = st.text_input("Scegli un Username").lower()
        password = st.text_input("Scegli una Password", type="password")
        password_confirm = st.text_input("Conferma Password", type="password")
        submitted = st.form_submit_button("Registrati")

        if submitted:
            users = load_users()
            if not username or not password:
                st.error("Username e password non possono essere vuoti.")
            elif username in users:
                st.error("Questo username è già stato preso.")
            elif password != password_confirm:
                st.error("Le password non coincidono.")
            else:
                users[username] = {"password": hash_password(password)}
                save_users(users)
                st.success("Registrazione completata! Ora puoi effettuare il login.")

def get_document_sources_from_vector_store():
    index_path = os.path.join(VECTOR_STORE_PATH, "index.faiss")
    if not os.path.exists(index_path):
        return []

    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=GOOGLE_API_KEY)
        vector_store = FAISS.load_local(VECTOR_STORE_PATH, embeddings, allow_dangerous_deserialization=True)
        
        all_metadata = [doc.metadata for doc in vector_store.docstore._dict.values()]
        sources = [meta.get('source', 'Sorgente Sconosciuta') for meta in all_metadata]
        
        unique_sources = sorted(list(set(sources)))
        return unique_sources
    except Exception as e:
        st.warning(f"Impossibile leggere le fonti dal vector store: {e}")
        return []

def render_main_app():
    st.title("💡 Consip Advisor")
    
    if "conversation" not in st.session_state:
        st.session_state.conversation = None

    if st.session_state.conversation is None and os.path.exists(os.path.join(VECTOR_STORE_PATH, "index.faiss")):
        try:
            embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=GOOGLE_API_KEY)
            vector_store = FAISS.load_local(VECTOR_STORE_PATH, embeddings, allow_dangerous_deserialization=True)
            st.session_state.conversation = get_conversational_chain(vector_store)
        except Exception as e: st.error(f"Impossibile caricare la base di conoscenza: {e}")

    with st.sidebar:
        st.header(f"Ciao, {st.session_state.current_user.capitalize()}")
        if st.button("Logout"):
            for key in list(st.session_state.keys()):
                del st.session_state[key]
            st.rerun()
        
        st.markdown("---")
        
        st.header("Le tue Conversazioni")
        if st.button("➕ Nuova Chat"):
            create_new_chat()
            save_active_chat(st.session_state.current_user)
        
        chats = st.session_state.user_data.get("chats", {})
        sorted_chats = sorted(chats.items(), key=lambda item: item[0], reverse=True)
        for chat_id, chat_data in sorted_chats:
            if st.button(chat_data["name"], key=chat_id, use_container_width=True):
                switch_chat(chat_id)
        
        st.markdown("---")
        st.header("Base di Conoscenza Condivisa")
        with st.expander("Documenti Caricati"):
            document_sources = get_document_sources_from_vector_store()
            if document_sources:
                for source in document_sources:
                    st.markdown(f"• `{source}`")
            else:
                st.info("Nessun documento è stato ancora processato.")
        
        allowed_types = ['pdf', 'docx', 'csv', 'xlsx', 'zip', '7z', 'pptx', 'aspx', 'txt']
        uploaded_files = st.file_uploader("Carica file o archivi (.zip, .7z) per le cartelle", accept_multiple_files=True, type=allowed_types)

        if st.button("Processa e Aggiungi"):
            if uploaded_files:
                with st.spinner("Elaborazione in corso..."):
                    docs = get_documents_with_detailed_metadata(uploaded_files)
                    if docs:
                        chunks = get_text_chunks(docs)
                        vs = get_vector_store(chunks)
                        if vs:
                            st.session_state.conversation = get_conversational_chain(vs)
                            st.success("Base di conoscenza aggiornata!")
                            st.rerun()
                    else: st.warning("Nessun documento valido trovato nei file caricati.")
            else: st.warning("Per favore, carica almeno un file.")
        
        st.markdown("---")
        st.header("Importa da Web")
        if st.button("Importa FAQ da AcquistinRetePA"):
            with st.spinner("Importazione dal web in corso..."):
                scraped_docs = scrape_faq_data()
                if scraped_docs:
                    chunks = get_text_chunks(scraped_docs)
                    vs = get_vector_store(chunks)
                    if vs:
                        st.session_state.conversation = get_conversational_chain(vs)
                        st.success("FAQ importate e base di conoscenza aggiornata!")
                        st.rerun()
                else:
                    st.error("Nessuna FAQ trovata o errore durante l'importazione.")
    
    for i, message in enumerate(st.session_state.get('messages', [])):
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
            if message["role"] == "assistant":
                citations = message.get("citations", [])
                if citations:
                    with st.expander("Mostra fonti consultate"):
                        for citation in citations: st.markdown(f"- {citation}")

                with st.popover("✏️ Fornisci un Feedback"):
                    with st.form(key=f"feedback_form_{i}"):
                        rating = st.radio("Valutazione:", ("Positivo 👍", "Negativo 👎"), horizontal=True)
                        comment = st.text_area("Commento per migliorare:", height=100)
                        if st.form_submit_button("Invia Feedback"):
                            save_feedback({
                                "timestamp": datetime.now().isoformat(), 
                                "username": st.session_state.current_user, 
                                "question": message.get("question", "N/D"), 
                                "response": message["content"], 
                                "rating": rating.split(" ")[0], 
                                "comment": comment
                            })
                            st.toast("Grazie! Feedback salvato.")
                            
    if user_question := st.chat_input("Fai una domanda sui documenti..."):
        if st.session_state.conversation is None:
            st.warning("La base di conoscenza è vuota. Per favore, carica dei documenti o importa le FAQ per iniziare a chattare.")
        else:
            active_chat_id = st.session_state.user_data.get("active_chat_id")
            if active_chat_id:
                user_message = {"role": "user", "content": user_question}
                st.session_state.user_data["chats"][active_chat_id]["messages"].append(user_message)
                st.session_state.messages = st.session_state.user_data["chats"][active_chat_id]["messages"]
                
                with st.chat_message("user"): st.markdown(user_question)
                
                with st.spinner("Consip Advisor sta pensando..."):
                    handle_user_input(user_question)
                
                save_active_chat(st.session_state.current_user)
                st.rerun()
            else:
                st.warning("Per favore, crea una 'Nuova Chat' per iniziare.")

if __name__ == "__main__":
    st.set_page_config(page_title="Consip Advisor", layout="wide")
    for path in [TEMP_FILES_PATH, VECTOR_STORE_PATH, USER_DATA_PATH]:
        if not os.path.exists(path): os.makedirs(path)

    if not st.session_state.get("logged_in"):
        show_login_or_register()
    else:
        if 'user_data' not in st.session_state:
            st.session_state.user_data = load_user_data(st.session_state.current_user)
            active_id = st.session_state.user_data.get("active_chat_id")
            
            if not active_id or not st.session_state.user_data["chats"]:
                create_new_chat()
            else:
                switch_chat(active_id)
        
        render_main_app()
